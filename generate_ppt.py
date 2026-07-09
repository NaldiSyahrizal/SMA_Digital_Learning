from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Pengembangan Aplikasi Digital Learning Terintegrasi Rekomendasi Prioritas SAW"
subtitle.text = "Nama Kelompok / Mahasiswa\nMata Kuliah / Program Studi"

# Slide 2: Rich Picture
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Rich Picture: Alur Sistem"
tf = body_shape.text_frame
tf.text = "Sistem Berjalan (Saat Ini):"
p = tf.add_paragraph()
p.text = "- Guru memberi tugas via grup WA/Classroom dasar."
p.level = 1
p = tf.add_paragraph()
p.text = "- Siswa sering bingung prioritas karena tugas menumpuk dan sering telat."
p.level = 1
p = tf.add_paragraph()
p.text = "Sistem Usulan (Aplikasi Kita):"
p.level = 0
p = tf.add_paragraph()
p.text = "- Terintegrasi dengan algoritma SAW (Simple Additive Weighting)."
p.level = 1
p = tf.add_paragraph()
p.text = "- Aplikasi merekomendasikan prioritas tugas cerdas untuk siswa."
p.level = 1

# Slide 3: UI and Mock-Up
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Rancangan UI dan Mock-Up"
tf = body_shape.text_frame
tf.text = "(Tempatkan Screenshot Aplikasi Di Sini)"
p = tf.add_paragraph()
p.text = "1. Beranda Siswa (Daftar Prioritas SAW)"
p.level = 1
p = tf.add_paragraph()
p.text = "2. Halaman Guru (Form Tambah Tugas & Kuis)"
p.level = 1
p = tf.add_paragraph()
p.text = "3. Halaman Admin (Kelola Master Data)"
p.level = 1

# Slide 4: Actors
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Aktor dan Hak Akses"
tf = body_shape.text_frame
tf.text = "1. Admin"
p = tf.add_paragraph()
p.text = "- Mengelola data master (Siswa, Guru, Mapel, Kelas)."
p.level = 1
p = tf.add_paragraph()
p.text = "2. Guru"
p.level = 0
p = tf.add_paragraph()
p.text = "- Membuat materi, tugas, dan menilai hasil ujian siswa."
p.level = 1
p = tf.add_paragraph()
p.text = "3. Siswa"
p.level = 0
p = tf.add_paragraph()
p.text = "- Mendapatkan rekomendasi prioritas tugas, melihat materi, dan mengerjakan kuis."
p.level = 1

# Slide 5: Features
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Fitur Utama Aplikasi"
tf = body_shape.text_frame
tf.text = "- Sistem Rekomendasi SAW (Tenggat Waktu, Bobot Tugas, Minat)."
p = tf.add_paragraph()
p.text = "- Penguncian Otomatis (Tugas ditutup setelah due date)."
p = tf.add_paragraph()
p.text = "- Kuis Pilihan Ganda dengan Auto-Grading."
p = tf.add_paragraph()
p.text = "- Notifikasi Real-time Publikasi dan Penilaian Tugas."

# Slide 6: Limitations
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Batasan Aplikasi"
tf = body_shape.text_frame
tf.text = "- Fokus pada evaluasi harian (PR, Proyek, Kuis), tidak mencakup UTS/UAS skala besar."
p = tf.add_paragraph()
p.text = "- Penilaian otomatis hanya untuk format soal pilihan ganda."
p = tf.add_paragraph()
p.text = "- Aplikasi difokuskan pada manajemen waktu siswa (Time Management)."

# Slide 7: Goals and Benefits
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Tujuan dan Manfaat"
tf = body_shape.text_frame
tf.text = "Tujuan:"
p = tf.add_paragraph()
p.text = "- Membangun platform e-learning yang bertindak sebagai asisten manajemen waktu."
p.level = 1
p = tf.add_paragraph()
p.text = "Manfaat:"
p.level = 0
p = tf.add_paragraph()
p.text = "- Siswa terhindar dari tugas menumpuk dan keterlambatan."
p.level = 1
p = tf.add_paragraph()
p.text = "- Guru mudah melacak kedisiplinan dan mempercepat proses grading."
p.level = 1

# Slide 8: Observation
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Bukti Observasi dan Wawancara"
tf = body_shape.text_frame
tf.text = "(Masukkan Foto Observasi Di Sini)"
p = tf.add_paragraph()
p.text = "Kutipan Wawancara Siswa:"
p.level = 0
p = tf.add_paragraph()
p.text = "'- Sering bingung memprioritaskan tugas saat deadlinenya berdekatan.'"
p.level = 1
p = tf.add_paragraph()
p.text = "Kutipan Wawancara Guru:"
p.level = 0
p = tf.add_paragraph()
p.text = "'- Butuh aplikasi yang membatasi pengumpulan tugas setelah lewat waktu.'"
p.level = 1

prs.save('Presentasi_PA_DigitalLearning.pptx')
print("PPTX generated successfully.")
